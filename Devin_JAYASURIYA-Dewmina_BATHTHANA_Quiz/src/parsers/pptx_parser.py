"""Parseur pour les documents PowerPoint (.pptx)."""

from pathlib import Path
from typing import List, Dict, Any
from pptx import Presentation

from .base_parser import BaseParser


class PptxParser(BaseParser):
    """Parseur pour les documents PowerPoint."""

    def parse(self) -> List[Dict[str, Any]]:
        """Parse le PowerPoint et retourne une liste de slides."""
        if not self.validate():
            raise FileNotFoundError(f"Le fichier {self.file_path} n'existe pas.")

        presentation = Presentation(self.file_path)
        slides = []

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

            slides.append({
                "slide": slide_num,
                "title": slide.shapes.title.text if slide.shapes.title else "",
                "content": "\n".join(slide_content),
                "text": "\n".join(slide_content)
            })

        return slides

    def extract_text(self) -> str:
        """Extract tout le texte du document."""
        if not self.validate():
            raise FileNotFoundError(f"Le fichier {self.file_path} n'existe pas.")

        presentation = Presentation(self.file_path)
        text = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text.strip())

        return "\n".join(text)
